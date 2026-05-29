import os
import sys
from datetime import datetime, timezone
from pathlib import Path
import collections.abc

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

from app.core.config import settings
from app.models.case import Case
from app.models.report_version import ReportVersion
from app.models.follow_up_entry import FollowUpEntry

# ─────────────────────────────────────────────────────────────────────────────
# ASSET PATHS
# ─────────────────────────────────────────────────────────────────────────────
_HERE = os.path.dirname(os.path.abspath(__file__))
_PROJECT_ROOT = os.path.abspath(os.path.join(_HERE, "../../.."))

# When running as a PyInstaller-frozen bundle, __file__ lives inside the
# temporary extraction directory and frontend/ is never present there.
# Assets are bundled via the `datas` key in backend.spec and land under
# sys._MEIPASS at the destination paths specified there.
def _asset(relative_path: str) -> str:
    """Resolve an asset path for both dev and frozen (PyInstaller) runs."""
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        return os.path.join(sys._MEIPASS, relative_path)
    
    # In Dev Mode, assets are located in the frontend source tree
    if relative_path.startswith("assets" + os.sep) or relative_path.startswith("assets/"):
        return os.path.join(_PROJECT_ROOT, "frontend", "src", relative_path)
        
    return os.path.join(_PROJECT_ROOT, relative_path)

LOGO_PATH      = _asset(os.path.join("assets", "logowhite.png"))
LOGO_DARK_PATH = _asset(os.path.join("assets", "logo.png"))

# ─────────────────────────────────────────────────────────────────────────────
# COLOURS
# ─────────────────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x16, 0x2C, 0x41)
SLATE = RGBColor(0x4F, 0x60, 0x6F)
MUTED = RGBColor(0x7C, 0x8A, 0x97)
FAINT = RGBColor(0xA7, 0xB4, 0xC0)
SURFACE_MUTED = RGBColor(0xF2, 0xF6, 0xFA)
BORDER_DEFAULT = RGBColor(0xD7, 0xE2, 0xEC)
BORDER_SUBTLE = RGBColor(0xE7, 0xEE, 0xF5)
BG_PRIMARY = RGBColor(0xF5, 0xF8, 0xFB)
BG_SECONDARY = RGBColor(0xED, 0xF3, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

C_STABLE = RGBColor(0x3E, 0x6B, 0x61)
C_EVOLVING = RGBColor(0xC5, 0x8A, 0x2B)
C_DECLINED = RGBColor(0x9B, 0x4A, 0x4A)
C_INFO = RGBColor(0x24, 0x4B, 0x73)

SEVERITY_COLORS = {
    "low": C_STABLE,
    "moderate": C_EVOLVING,
    "high": C_DECLINED,
    "critical": C_DECLINED,
}

URGENCY_COLORS = {
    "routine": MUTED,
    "urgent": C_EVOLVING,
    "emergent": C_DECLINED,
}

class PPTXExportService:
    """Service to generate PowerPoint presentations for clinical reports."""

    def __init__(self):
        self.output_dir = settings.pdf_output_dir
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir, exist_ok=True)

    def _preprocess_citations(self, report: ReportVersion) -> ReportVersion:
        """Recursively replace {{cite:cN}} with [N] in all text fields of the report."""
        import re, copy, collections.abc
        citations = report.citations or []
        def repl(m):
            idx = int(m.group(1))
            return f"[{idx}]" if 1 <= idx <= len(citations) else m.group(0)
        def walk(obj):
            if isinstance(obj, dict):
                return {k: walk(v) for k, v in obj.items()}
            elif isinstance(obj, list):
                return [walk(i) for i in obj]
            elif isinstance(obj, str):
                return re.sub(r'\{\{cite:c?(\d+)\}\}', repl, obj)
            else:
                return obj
        # Shallow copy to avoid mutating original
        new_report = copy.copy(report)
        # Process all relevant fields
        new_report.summary = walk(report.summary)
        new_report.differentials = walk(report.differentials)
        new_report.supporting_findings = walk(report.supporting_findings)
        new_report.contradictory_findings = walk(report.contradictory_findings)
        new_report.next_steps = walk(report.next_steps)
        new_report.missing_information = walk(report.missing_information)
        new_report.citations = walk(report.citations)  # optional, but harmless
        return new_report

    def generate_report_pptx(self, report: ReportVersion, case: Case, follow_ups: list[FollowUpEntry]) -> str:
        report = self._preprocess_citations(report)   
        filename = f"report_{report.id}.pptx"
        output_path = os.path.join(self.output_dir, filename)

        prs = Presentation()
        # Set slide size to standard widescreen (13.33" × 7.5")
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Blank slide layout is typically layout 6 in the default template
        blank_layout = prs.slide_layouts[6]

        total_slides_placeholder = [] # We'll fill this with refs to the slides so we can add footers later

        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 1: COVER
        # ─────────────────────────────────────────────────────────────────────
        slide_cover = prs.slides.add_slide(blank_layout)
        # Background
        bg = slide_cover.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()

        # Logo
        if os.path.exists(LOGO_PATH):
            try:
                # Top-left, ~35% of slide width
                logo_width = prs.slide_width * 0.14
                slide_cover.shapes.add_picture(
                    LOGO_PATH, Inches(0.6), Inches(0.6), logo_width
                )
            except Exception:
                pass

        # Title
        tx_title = slide_cover.shapes.add_textbox(
            Inches(1), prs.slide_height - Inches(3.5), prs.slide_width - Inches(2), Inches(1)
        )
        p_title = tx_title.text_frame.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        run_title = p_title.add_run()
        run_title.text = "CLINICAL REASONING REPORT"
        run_title.font.name = "TeX Gyre Termes"
        run_title.font.size = Pt(38)
        run_title.font.bold = True
        run_title.font.color.rgb = WHITE

        # Subtitle 1 (Case Title)
        tx_sub1 = slide_cover.shapes.add_textbox(
            Inches(1), prs.slide_height - Inches(2.3), prs.slide_width - Inches(2), Inches(0.5)
        )
        p_sub1 = tx_sub1.text_frame.paragraphs[0]
        p_sub1.alignment = PP_ALIGN.CENTER
        run_sub1 = p_sub1.add_run()
        run_sub1.text = case.title or "Untitled Case"
        run_sub1.font.name = "TeX Gyre Termes"
        run_sub1.font.size = Pt(18)
        run_sub1.font.color.rgb = FAINT

        # Subtitle 2 (Version & Date)
        gen_date = report.created_at.strftime("%Y-%m-%d %H:%M UTC") if report.created_at else datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        tx_sub2 = slide_cover.shapes.add_textbox(
            Inches(1), prs.slide_height - Inches(1.8), prs.slide_width - Inches(2), Inches(0.5)
        )
        p_sub2 = tx_sub2.text_frame.paragraphs[0]
        p_sub2.alignment = PP_ALIGN.CENTER
        run_sub2 = p_sub2.add_run()
        run_sub2.text = f"Version {report.version_number or 1}  ·  {gen_date}"
        run_sub2.font.name = "TeX Gyre Termes"
        run_sub2.font.size = Pt(13)
        run_sub2.font.color.rgb = FAINT

        # Severity Badge
        summary_data = report.summary or {}
        if isinstance(summary_data, dict):
            severity = summary_data.get("severity")
            if severity:
                sev_color = SEVERITY_COLORS.get(severity.lower(), C_EVOLVING)
                badge_w = Inches(1.5)
                badge_h = Inches(0.4)
                badge_l = prs.slide_width - badge_w - Inches(0.5)
                badge_t = prs.slide_height - badge_h - Inches(0.5)
                
                badge_shape = slide_cover.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, badge_l, badge_t, badge_w, badge_h
                )
                badge_shape.fill.solid()
                badge_shape.fill.fore_color.rgb = sev_color
                badge_shape.line.fill.background()
                
                bp = badge_shape.text_frame.paragraphs[0]
                bp.alignment = PP_ALIGN.CENTER
                br = bp.add_run()
                br.text = severity.upper()
                br.font.name = "TeX Gyre Termes"
                br.font.size = Pt(10)
                br.font.bold = True
                br.font.color.rgb = WHITE


        # ─────────────────────────────────────────────────────────────────────
        # HELPER: ADD FOOTER
        # ─────────────────────────────────────────────────────────────────────
        def _add_footer(slide, page_num, total_pages):
            # Footer bar
            f_h = Pt(20)
            f_t = prs.slide_height - f_h
            f_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, f_t, prs.slide_width, f_h
            )
            f_shape.fill.solid()
            f_shape.fill.fore_color.rgb = BG_PRIMARY
            f_shape.line.color.rgb = BORDER_SUBTLE
            f_shape.line.width = Pt(0.5)

            # Left text
            tx_left = slide.shapes.add_textbox(Inches(0.5), f_t, prs.slide_width / 2, f_h)
            tx_left.text_frame.word_wrap = False
            tx_left.text_frame.margin_top = Pt(2)
            pl = tx_left.text_frame.paragraphs[0]
            pl.alignment = PP_ALIGN.LEFT
            rl = pl.add_run()
            rl.text = "Aletheia Clinical Workstation  ·  Educational / research use only"
            rl.font.name = "TeX Gyre Termes"
            rl.font.size = Pt(7.5)
            rl.font.color.rgb = FAINT

            # Right text
            tx_right = slide.shapes.add_textbox(prs.slide_width / 2, f_t, prs.slide_width / 2 - Inches(0.5), f_h)
            tx_right.text_frame.word_wrap = False
            tx_right.text_frame.margin_top = Pt(2)
            pr = tx_right.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.RIGHT
            rr = pr.add_run()
            rr.text = f"Page {page_num} of {total_pages}"
            rr.font.name = "TeX Gyre Termes"
            rr.font.size = Pt(7.5)
            rr.font.color.rgb = FAINT

        # ─────────────────────────────────────────────────────────────────────
        # HELPER: HEADER FOR CONTENT SLIDES
        # ─────────────────────────────────────────────────────────────────────
        def _add_content_header(slide, title_text):
            tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), prs.slide_width - Inches(1.6), Inches(0.6))
            p_title = tx_title.text_frame.paragraphs[0]
            r_title = p_title.add_run()
            r_title.text = title_text
            r_title.font.name = "TeX Gyre Termes"
            r_title.font.size = Pt(22)
            r_title.font.bold = True
            r_title.font.color.rgb = NAVY

            # Horizontal rule
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(0.8), Inches(1.1), prs.slide_width - Inches(0.8), Inches(1.1)
            )
            line.line.color.rgb = BORDER_DEFAULT
            line.line.width = Pt(1)

        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 2: CLINICAL SUMMARY
        # ─────────────────────────────────────────────────────────────────────
        if isinstance(summary_data, dict) and summary_data.get("overview"):
            slide_summary = prs.slides.add_slide(blank_layout)
            total_slides_placeholder.append(slide_summary)
            _add_content_header(slide_summary, "Clinical Summary")
            
            # Label
            tx_label = slide_summary.shapes.add_textbox(prs.slide_width - Inches(2), Inches(0.2), Inches(1.5), Inches(0.4))
            p_label = tx_label.text_frame.paragraphs[0]
            p_label.alignment = PP_ALIGN.RIGHT
            r_label = p_label.add_run()
            r_label.text = "ALETHEIA"
            r_label.font.name = "TeX Gyre Termes"
            r_label.font.size = Pt(8)
            r_label.font.color.rgb = MUTED

            # Body text
            tx_body = slide_summary.shapes.add_textbox(Inches(0.8), Inches(1.5), prs.slide_width - Inches(1.6), prs.slide_height - Inches(2.5))
            tx_body.text_frame.word_wrap = True
            p_body = tx_body.text_frame.paragraphs[0]
            p_body.alignment = PP_ALIGN.LEFT
            r_body = p_body.add_run()
            r_body.text = summary_data.get("overview")
            r_body.font.name = "TeX Gyre Termes"
            r_body.font.size = Pt(12)
            r_body.font.color.rgb = NAVY


        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 3: ASSESSMENT & DIFFERENTIALS
        # ─────────────────────────────────────────────────────────────────────
        differentials = report.differentials or []
        if differentials:
            # Sort descending by confidence
            differentials = sorted(differentials, key=lambda x: x.get("confidence", 0), reverse=True)
            
            def _create_diff_slide():
                slide = prs.slides.add_slide(blank_layout)
                total_slides_placeholder.append(slide)
                _add_content_header(slide, "Assessment & Differential Diagnoses")
                return slide

            slide_diff = _create_diff_slide()
            current_y = Inches(1.5)
            items_on_slide = 0

            for dx in differentials:
                if items_on_slide >= 4:
                    slide_diff = _create_diff_slide()
                    current_y = Inches(1.5)
                    items_on_slide = 0

                name = dx.get("diagnosis", "Unknown Diagnosis")
                conf = dx.get("confidence", 0)
                reasoning = dx.get("reasoning", "")

                conf_color = C_INFO if conf >= 0.75 else (RGBColor(0x5D, 0x78, 0x92) if conf >= 0.5 else RGBColor(0xA5, 0xB4, 0xC2))

                # Name
                tx_name = slide_diff.shapes.add_textbox(Inches(0.8), current_y, Inches(5), Inches(0.4))
                pn = tx_name.text_frame.paragraphs[0]
                rn = pn.add_run()
                rn.text = name
                rn.font.name = "TeX Gyre Termes"
                rn.font.size = Pt(16)
                rn.font.bold = True
                rn.font.color.rgb = NAVY

                # Confidence %
                tx_conf = slide_diff.shapes.add_textbox(prs.slide_width - Inches(2), current_y, Inches(1.2), Inches(0.4))
                pc = tx_conf.text_frame.paragraphs[0]
                pc.alignment = PP_ALIGN.RIGHT
                rc = pc.add_run()
                rc.text = f"{conf*100:.0f}%"
                rc.font.name = "TeX Gyre Termes"
                rc.font.size = Pt(15)
                rc.font.bold = True
                rc.font.color.rgb = conf_color

                current_y += Inches(0.35)

                # Progress bar
                bar_w = prs.slide_width - Inches(1.6)
                fill_w = bar_w * min(max(conf, 0), 1)
                
                if fill_w > 0:
                    bar_fill = slide_diff.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), current_y, fill_w, Pt(6))
                    bar_fill.fill.solid()
                    bar_fill.fill.fore_color.rgb = conf_color
                    bar_fill.line.fill.background()

                if bar_w - fill_w > 0:
                    bar_empty = slide_diff.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8) + fill_w, current_y, bar_w - fill_w, Pt(6))
                    bar_empty.fill.solid()
                    bar_empty.fill.fore_color.rgb = BG_SECONDARY
                    bar_empty.line.fill.background()

                current_y += Pt(6) + Inches(0.1)

                # Reasoning
                if reasoning:
                    tx_reason = slide_diff.shapes.add_textbox(Inches(0.8), current_y, bar_w, Inches(0.5))
                    tx_reason.text_frame.word_wrap = True
                    pr = tx_reason.text_frame.paragraphs[0]
                    rr = pr.add_run()
                    rr.text = reasoning
                    rr.font.name = "TeX Gyre Termes"
                    rr.font.size = Pt(11.5)
                    rr.font.italic = True
                    rr.font.color.rgb = SLATE
                    current_y += tx_reason.height

                current_y += Pt(8)
                items_on_slide += 1


        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 4: SUPPORTING & CONTRADICTORY FINDINGS
        # ─────────────────────────────────────────────────────────────────────
        supp_findings = report.supporting_findings or []
        cont_findings = report.contradictory_findings or []
        
        if supp_findings or cont_findings:
            def _create_findings_slide():
                slide = prs.slides.add_slide(blank_layout)
                total_slides_placeholder.append(slide)
                _add_content_header(slide, "Supporting & Contradictory Findings")
                
                # Center divider
                div_x = prs.slide_width / 2
                div = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, div_x, Inches(1.3), div_x, prs.slide_height - Inches(0.8))
                div.line.color.rgb = BORDER_SUBTLE
                div.line.width = Pt(0.5)

                # Headers
                # Left
                tx_lh = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), div_x - Inches(1.6), Inches(0.4))
                pl = tx_lh.text_frame.paragraphs[0]
                rl = pl.add_run()
                rl.text = "Supporting Findings"
                rl.font.name = "TeX Gyre Termes"
                rl.font.size = Pt(13)
                rl.font.bold = True
                rl.font.color.rgb = C_STABLE

                # Right
                tx_rh = slide.shapes.add_textbox(div_x + Inches(0.8), Inches(1.3), div_x - Inches(1.6), Inches(0.4))
                pr = tx_rh.text_frame.paragraphs[0]
                rr = pr.add_run()
                rr.text = "Contradictory Findings"
                rr.font.name = "TeX Gyre Termes"
                rr.font.size = Pt(13)
                rr.font.bold = True
                rr.font.color.rgb = C_EVOLVING
                
                return slide

            slide_find = _create_findings_slide()
            y_left = Inches(1.8)
            y_right = Inches(1.8)
            max_y = prs.slide_height - Inches(1)
            
            def _add_finding(slide, is_left, title, text, strength):
                nonlocal y_left, y_right, slide_find
                current_y = y_left if is_left else y_right
                x_pos = Inches(0.8) if is_left else (prs.slide_width / 2 + Inches(0.8))
                w_pos = prs.slide_width / 2 - Inches(1.6)

                if current_y > max_y - Inches(0.8):
                    slide_find = _create_findings_slide()
                    y_left = Inches(1.8)
                    y_right = Inches(1.8)
                    current_y = y_left if is_left else y_right

                tx = slide_find.shapes.add_textbox(x_pos, current_y, w_pos, Inches(1))
                tx.text_frame.word_wrap = True
                
                p1 = tx.text_frame.paragraphs[0]
                r1 = p1.add_run()
                r1.text = title + "\n"
                r1.font.name = "TeX Gyre Termes"
                r1.font.size = Pt(10)
                r1.font.bold = True
                r1.font.color.rgb = NAVY

                p2 = tx.text_frame.add_paragraph()
                r2 = p2.add_run()
                r2.text = text + "\n"
                r2.font.name = "TeX Gyre Termes"
                r2.font.size = Pt(9)
                r2.font.italic = True
                r2.font.color.rgb = SLATE

                if strength:
                    str_color = C_STABLE if strength.lower() == "strong" else (C_EVOLVING if strength.lower() == "moderate" else MUTED)
                    p3 = tx.text_frame.add_paragraph()
                    r3 = p3.add_run()
                    r3.text = strength.upper()
                    r3.font.name = "TeX Gyre Termes"
                    r3.font.size = Pt(8)
                    r3.font.color.rgb = str_color
                
                # Approximate height bump
                added_h = Inches(0.6) + min(len(text)//40, 5) * Inches(0.15)
                if is_left:
                    y_left += added_h
                else:
                    y_right += added_h

            for sf in supp_findings:
                _add_finding(slide_find, True, sf.get("finding", "Unknown"), sf.get("explanation", ""), sf.get("significance", sf.get("strength", "")))
            
            for cf in cont_findings:
                _add_finding(slide_find, False, cf.get("finding", "Unknown"), cf.get("explanation", ""), cf.get("significance", cf.get("strength", "")))

        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 5: SUGGESTED NEXT STEPS
        # ─────────────────────────────────────────────────────────────────────
        next_steps = report.next_steps or []
        if next_steps:
            def _create_steps_slide():
                slide = prs.slides.add_slide(blank_layout)
                total_slides_placeholder.append(slide)
                _add_content_header(slide, "Suggested Next Steps")
                return slide

            slide_steps = _create_steps_slide()
            current_y = Inches(1.5)

            for idx, ns in enumerate(next_steps, 1):
                if current_y > prs.slide_height - Inches(1.5):
                    slide_steps = _create_steps_slide()
                    current_y = Inches(1.5)

                title = ns.get("title", ns.get("action", "Unknown Action"))
                urgency = ns.get("urgency", "routine")
                urg_color = URGENCY_COLORS.get(urgency.lower(), MUTED)
                rationale = ns.get("rationale", "")
                risks = ns.get("risks_of_delay", "")

                tx = slide_steps.shapes.add_textbox(Inches(0.8), current_y, prs.slide_width - Inches(1.6), Inches(1))
                tx.text_frame.word_wrap = True

                p1 = tx.text_frame.paragraphs[0]
                r1 = p1.add_run()
                r1.text = f"{idx}. {title}  "
                r1.font.name = "TeX Gyre Termes"
                r1.font.size = Pt(18)
                r1.font.bold = True
                r1.font.color.rgb = NAVY

                r_urg = p1.add_run()
                r_urg.text = urgency.upper()
                r_urg.font.name = "TeX Gyre Termes"
                r_urg.font.size = Pt(14)
                r_urg.font.bold = True
                r_urg.font.color.rgb = urg_color

                if rationale:
                    p2 = tx.text_frame.add_paragraph()
                    r2 = p2.add_run()
                    r2.text = rationale
                    r2.font.name = "TeX Gyre Termes"
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = SLATE
                
                if risks:
                    p3 = tx.text_frame.add_paragraph()
                    r3 = p3.add_run()
                    r3.text = f"Risk of delay: {risks}"
                    r3.font.name = "TeX Gyre Termes"
                    r3.font.size = Pt(10)
                    r3.font.italic = True
                    r3.font.color.rgb = C_DECLINED

                current_y += Inches(0.6) + min((len(rationale) + len(risks))//80, 4) * Inches(0.15)


        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 6: MISSING INFORMATION
        # ─────────────────────────────────────────────────────────────────────
        missing_info = report.missing_information or []
        if missing_info:
            def _create_missing_slide():
                slide = prs.slides.add_slide(blank_layout)
                total_slides_placeholder.append(slide)
                _add_content_header(slide, "Information Gaps")
                return slide

            slide_missing = _create_missing_slide()
            current_y = Inches(1.5)

            for mi in missing_info:
                if current_y > prs.slide_height - Inches(1.5):
                    slide_missing = _create_missing_slide()
                    current_y = Inches(1.5)

                item = mi.get("item", mi.get("information", "Unknown Item"))
                cat = mi.get("category", "")
                why = mi.get("why_it_matters", mi.get("reason", ""))
                impact = mi.get("impact_on_assessment", "")

                tx = slide_missing.shapes.add_textbox(Inches(0.8), current_y, prs.slide_width - Inches(1.6), Inches(1))
                tx.text_frame.word_wrap = True

                p1 = tx.text_frame.paragraphs[0]
                r1 = p1.add_run()
                r1.text = item + "  "
                r1.font.name = "TeX Gyre Termes"
                r1.font.size = Pt(18)
                r1.font.bold = True
                r1.font.color.rgb = NAVY

                if cat:
                    rc = p1.add_run()
                    rc.text = f"[{cat.upper()}]"
                    rc.font.name = "TeX Gyre Termes"
                    rc.font.size = Pt(10.5)
                    rc.font.color.rgb = MUTED

                if why:
                    p2 = tx.text_frame.add_paragraph()
                    r2 = p2.add_run()
                    r2.text = why
                    r2.font.name = "TeX Gyre Termes"
                    r2.font.size = Pt(13)
                    r2.font.color.rgb = SLATE
                
                if impact:
                    p3 = tx.text_frame.add_paragraph()
                    r3 = p3.add_run()
                    r3.text = f"Impact: {impact}"
                    r3.font.name = "TeX Gyre Termes"
                    r3.font.size = Pt(13)
                    r3.font.italic = True
                    r3.font.color.rgb = SLATE

                current_y += Inches(1.2) + Inches(0.18)


        # ─────────────────────────────────────────────────────────────────────
        # SLIDE 7: REFERENCES
        # ─────────────────────────────────────────────────────────────────────
        citations = report.citations or []
        if citations:
            REF_LEFT = Inches(0.8)
            REF_RIGHT = Inches(0.8)
            REF_TOP = Inches(1.45)
            REF_BOTTOM_GUARD = Inches(0.75)

            BODY_SIZE = Pt(11)
            INDEX_SIZE = Pt(11)
            EV_SIZE = Pt(9)

            def _create_ref_slide():
                slide = prs.slides.add_slide(blank_layout)
                total_slides_placeholder.append(slide)
                _add_content_header(slide, "References")
                return slide

            def _estimate_ref_height(text: str, ev_level: str = "") -> Inches:
                # Rough but reliable enough for slide pagination.
                # Longer refs get more height so they do not collide.
                est_chars = len(text) + (len(ev_level) + 4 if ev_level else 0)
                lines = max(1, (est_chars // 115) + 1)
                return Inches(0.24 * lines + 0.10)

            slide_ref = _create_ref_slide()
            current_y = REF_TOP
            refs_on_slide = 0

            for idx, cit in enumerate(citations, 1):
                van = cit.get("vancouverString", "")
                if not van:
                    title = cit.get("title", cit.get("label", "Untitled"))
                    authors = cit.get("authors", [])
                    journal = cit.get("journal", "")
                    year = cit.get("year", "")
                    pmid = cit.get("pmid", "")

                    auth_str = ""
                    if authors:
                        listed = authors[:3]
                        suffix = " et al." if len(authors) > 3 else ""
                        auth_str = ", ".join(listed) + suffix + ". "

                    van = f"{auth_str}{title}. {journal}. {year}."
                    if pmid:
                        van += f" PMID {pmid}."

                ev_level = cit.get("evidenceLevel", "")
                est_h = _estimate_ref_height(van, ev_level)

                # Start a new slide if this one would overflow or if too many refs are already on it.
                if refs_on_slide >= 6 or current_y + est_h > prs.slide_height - REF_BOTTOM_GUARD:
                    slide_ref = _create_ref_slide()
                    current_y = REF_TOP
                    refs_on_slide = 0

                tx = slide_ref.shapes.add_textbox(
                    REF_LEFT,
                    current_y,
                    prs.slide_width - REF_LEFT - REF_RIGHT,
                    est_h,
                )
                tf = tx.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.space_before = Pt(0)
                p.space_after = Pt(8)

                rn = p.add_run()
                rn.text = f"[{idx}] "
                rn.font.name = "TeX Gyre Termes"
                rn.font.size = INDEX_SIZE
                rn.font.bold = True
                rn.font.color.rgb = C_INFO

                rt = p.add_run()
                rt.text = van + " "
                rt.font.name = "TeX Gyre Termes"
                rt.font.size = BODY_SIZE
                rt.font.color.rgb = SLATE

                if ev_level:
                    re = p.add_run()
                    re.text = f"[{ev_level.upper()}]"
                    re.font.name = "TeX Gyre Termes"
                    re.font.size = EV_SIZE
                    re.font.italic = True
                    re.font.color.rgb = MUTED

                current_y += est_h + Inches(0.08)
                refs_on_slide += 1

        # ─────────────────────────────────────────────────────────────────────
        # ADD FOOTERS TO ALL CONTENT SLIDES
        # ─────────────────────────────────────────────────────────────────────
        # slide 1 is cover, total_slides_placeholder has slide 2+
        # cover is page 1, so total = len(total_slides_placeholder) + 1
        total_pages = len(total_slides_placeholder) + 1
        for idx, slide in enumerate(total_slides_placeholder, 2):
            _add_footer(slide, idx, total_pages)


        # Save
        try:
            prs.save(output_path)
        except Exception as e:
            raise Exception(f"Failed to save PPTX: {str(e)}")

        return output_path
